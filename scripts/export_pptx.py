#!/usr/bin/env python3
"""Export CyberODM executive report to PowerPoint (.pptx).

Requires:
  pip install python-pptx pillow
Usage:
  python scripts/export_pptx.py --title "CyberODM Exec Update Q2 2025" --out outputs/CyberODM-Exec-Report.pptx
"""
import argparse
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide

def add_table_slide(prs, title, df):
    slide_layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    rows, cols = df.shape[0]+1, df.shape[1]
    x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(9), Inches(4.5)
    table = slide.shapes.add_table(rows, cols, x, y, cx, cy).table
    for j, col in enumerate(df.columns):
        table.cell(0, j).text = str(col)
    for i in range(df.shape[0]):
        for j in range(cols):
            table.cell(i+1, j).text = str(df.iloc[i, j])
    return slide

def add_picture_slide(prs, title, image_path):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.shapes.add_picture(str(image_path), Inches(0.5), Inches(1.5), height=Inches(4.5))
    return slide

def main():
    p = argparse.ArgumentParser()
    p.add_argument("--base", default=".")
    p.add_argument("--title", default="CyberODM Executive Update")
    p.add_argument("--subtitle", default="Outcome‑Driven Metrics")
    p.add_argument("--out", default="outputs/CyberODM-Exec-Report.pptx")
    args = p.parse_args()

    base = Path(args.base)
    outdir = base / "outputs"
    outdir.mkdir(exist_ok=True)

    import pandas as pd
    latest = pd.read_csv(outdir / "outcome_scores_latest.csv")
    val = pd.read_csv(outdir / "value_realization.csv")

    prs = Presentation()
    add_title_slide(prs, args.title, args.subtitle)

    # Snapshot
    cols = ["outcome_id","name","owner","outcome_score","rag","okr_progress_pct","value_realization_usd"]
    add_table_slide(prs, "Snapshot (Latest)", latest[cols])

    # Charts
    for title, img in [
        ("Outcome Scores", outdir / "outcome_scores_latest.png"),
        ("Value Realization Trend", outdir / "value_trend.png"),
        ("OKR Progress (Latest)", outdir / "okr_latest.png"),
    ]:
        if img.exists():
            add_picture_slide(prs, title, img)

    # SLA Overlay slide (if exists)
    sla = outdir / 'sla_overlay.png'
    if sla.exists():
        add_picture_slide(prs, 'SLA/SLO Overlay', sla)

    # Top 5 Actions slide (derive from lowest outcomes + initiatives)
    import pandas as pd
    latest = pd.read_csv(outdir / 'outcome_scores_latest.csv')
    inits = pd.read_csv(outdir / 'initiatives_snapshot.csv')
    gaps = latest.sort_values('outcome_score').head(5)
    actions = []
    for _, row in gaps.iterrows():
        out_id = row['outcome_id']
        cand = inits[inits['outcome_id'] == out_id]
        action = cand['name'].iloc[0] if not cand.empty else 'Accelerate improvement plan'
        owner = row['owner']
        actions.append([row['name'], owner, f"Focus: {action}", f"Score={row['outcome_score']:.2f}"])
    import pandas as pd
    actdf = pd.DataFrame(actions, columns=['Outcome','Owner','Action','Context'])
    add_table_slide(prs, 'Top 5 Actions (Next Quarter)', actdf)

    prs.save(str(base / args.out))
    print("Wrote:", base / args.out)

if __name__ == "__main__":
    main()
